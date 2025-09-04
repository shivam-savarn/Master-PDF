from flask import Flask, request, jsonify, send_file, after_this_request
from werkzeug.utils import secure_filename
import os
import PyPDF2
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
import tempfile
import zipfile
from datetime import datetime
import io
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("PyMuPDF not available, using basic compression")

app = Flask(__name__, static_folder='.', static_url_path='')
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

ALLOWED_EXTENSIONS = {
    'pdf': ['pdf'],
    'image': ['jpg', 'jpeg', 'png', 'gif', 'bmp']
}

def allowed_file(filename, file_type):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS[file_type]

def send_file_and_cleanup(file_path, filename):
    @after_this_request
    def remove_file(response):
        try:
            os.unlink(file_path)
        except Exception:
            pass
        return response
    
    return send_file(file_path, as_attachment=True, download_name=filename)

@app.route('/')
def index():
    return app.send_static_file('index.html')

@app.route('/<path:filename>')
def static_files(filename):
    return app.send_static_file(filename)

@app.route('/merge-pdf', methods=['POST'])
def merge_pdf():
    try:
        files = request.files.getlist('files')
        if len(files) < 2:
            return jsonify({'error': 'At least 2 PDF files required'}), 400
        
        for file in files:
            if not allowed_file(file.filename, 'pdf'):
                return jsonify({'error': f'Invalid file: {file.filename}'}), 400
        
        merger = PyPDF2.PdfMerger()
        
        for file in files:
            merger.append(file)
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'merged_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        with open(output_path, 'wb') as output_file:
            merger.write(output_file)
        
        merger.close()
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/jpg-to-pdf', methods=['POST'])
def jpg_to_pdf():
    try:
        files = request.files.getlist('files')
        page_size = request.form.get('pageSize', 'A4')
        orientation = request.form.get('orientation', 'portrait')
        
        if not files or len(files) == 0:
            return jsonify({'error': 'No files provided'}), 400
        
        for file in files:
            if not file.filename or not allowed_file(file.filename, 'image'):
                return jsonify({'error': f'Invalid file: {file.filename}'}), 400
        
        if page_size == 'A4':
            page_dims = A4
        elif page_size == 'Letter':
            page_dims = letter
        else:
            page_dims = A4
        
        if orientation == 'landscape':
            page_dims = (page_dims[1], page_dims[0])
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'images_to_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        c = canvas.Canvas(output_path, pagesize=page_dims)
        page_width, page_height = page_dims
        
        for i, file in enumerate(files):
            try:
                temp_input = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1])
                file.save(temp_input.name)
                temp_input.close()
                
                img = Image.open(temp_input.name)
                
                if img.mode != 'RGB':
                    if img.mode in ('RGBA', 'LA'):
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'RGBA':
                            background.paste(img, mask=img.split()[-1])
                        else:
                            background.paste(img)
                        img = background
                    else:
                        img = img.convert('RGB')
                
                img_width, img_height = img.size
                
                scale_x = (page_width - 40) / img_width
                scale_y = (page_height - 40) / img_height
                scale = min(scale_x, scale_y)
                
                new_width = img_width * scale
                new_height = img_height * scale
                
                x = (page_width - new_width) / 2
                y = (page_height - new_height) / 2
                
                temp_output = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
                img.save(temp_output.name, 'JPEG', quality=95)
                temp_output.close()
                
                c.drawImage(temp_output.name, x, y, new_width, new_height)
                c.showPage()
                
                os.unlink(temp_input.name)
                os.unlink(temp_output.name)
                    
            except Exception as img_error:
                return jsonify({'error': f'Error processing image {file.filename}: {str(img_error)}'}), 500
        
        c.save()
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'JPG to PDF conversion failed: {str(e)}'}), 500

@app.route('/pdf-to-ppt', methods=['POST'])
def pdf_to_ppt():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'pdf_to_ppt_{timestamp}.pptx'
        output_path = tempfile.mktemp(suffix='.pptx')
        
        convert_pdf_to_pptx(temp_input.name, output_path)
        
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'PDF to PowerPoint conversion failed: {str(e)}'}), 500

def convert_pdf_to_pptx(pdf_path, pptx_path):
    try:
        from pptx import Presentation
        
        if not PYMUPDF_AVAILABLE:
            raise Exception("PyMuPDF is required for PDF to PowerPoint conversion")
        
        doc_pdf = fitz.open(pdf_path)
        prs = Presentation()
        
        for page_num in range(len(doc_pdf)):
            try:
                page = doc_pdf[page_num]
                
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title_placeholder = slide.shapes.title
                content_placeholder = slide.placeholders[1]
                
                text_dict = page.get_text("dict")
                
                all_text = []
                for block in text_dict["blocks"]:
                    if "lines" in block:
                        block_text = ""
                        for line in block["lines"]:
                            line_text = ""
                            for span in line["spans"]:
                                line_text += span["text"]
                            if line_text.strip():
                                block_text += line_text.strip() + " "
                        
                        if block_text.strip():
                            all_text.append(block_text.strip())
                
                if all_text:
                    title_placeholder.text = f"Page {page_num + 1}"
                    content_placeholder.text = "\n".join(all_text[:10])
                else:
                    title_placeholder.text = f"Page {page_num + 1}"
                    content_placeholder.text = "No text content found on this page."
                    
            except Exception as page_error:
                print(f"Error processing page {page_num}: {str(page_error)}")
                continue
        
        doc_pdf.close()
        prs.save(pptx_path)
        
    except Exception as e:
        raise Exception(f"PDF to PowerPoint conversion error: {str(e)}")

@app.route('/add-page-numbers', methods=['POST'])
def add_page_numbers():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        position = request.form.get('position', 'bottom-right')
        start_page = int(request.form.get('startPage', 1))
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'numbered_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        add_page_numbers_to_pdf(temp_input.name, output_path, position, start_page)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Add page numbers failed: {str(e)}'}), 500

def add_page_numbers_to_pdf(input_path, output_path, position, start_page):
    reader = PyPDF2.PdfReader(input_path)
    writer = PyPDF2.PdfWriter()
    
    for i, page in enumerate(reader.pages):
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        
        if position == 'top-left':
            x, y = 50, 750
        elif position == 'top-right':
            x, y = 500, 750
        elif position == 'bottom-left':
            x, y = 50, 50
        else:
            x, y = 500, 50
        
        page_num = i + start_page
        can.drawString(x, y, str(page_num))
        can.save()
        
        packet.seek(0)
        overlay = PyPDF2.PdfReader(packet)
        page.merge_page(overlay.pages[0])
        writer.add_page(page)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

@app.route('/add-watermark', methods=['POST'])
def add_watermark():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        watermark_text = request.form.get('watermarkText', 'WATERMARK')
        opacity = float(request.form.get('opacity', 0.3))
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'watermarked_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        add_watermark_to_pdf(temp_input.name, output_path, watermark_text, opacity)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Add watermark failed: {str(e)}'}), 500

def add_watermark_to_pdf(input_path, output_path, watermark_text, opacity):
    reader = PyPDF2.PdfReader(input_path)
    writer = PyPDF2.PdfWriter()
    
    for page in reader.pages:
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        
        can.setFillAlpha(opacity)
        can.setFont("Helvetica-Bold", 50)
        
        can.saveState()
        can.translate(300, 400)
        can.rotate(45)
        
        text_width = can.stringWidth(watermark_text, "Helvetica-Bold", 50)
        can.drawString(-text_width/2, 0, watermark_text)
        
        can.restoreState()
        can.save()
        
        packet.seek(0)
        overlay = PyPDF2.PdfReader(packet)
        page.merge_page(overlay.pages[0])
        writer.add_page(page)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

@app.route('/crop-pdf', methods=['POST'])
def crop_pdf():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        top = float(request.form.get('top', 0)) / 100
        bottom = float(request.form.get('bottom', 0)) / 100
        left = float(request.form.get('left', 0)) / 100
        right = float(request.form.get('right', 0)) / 100
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'cropped_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        crop_pdf_pages(temp_input.name, output_path, top, bottom, left, right)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Crop PDF failed: {str(e)}'}), 500

def crop_pdf_pages(input_path, output_path, top, bottom, left, right):
    reader = PyPDF2.PdfReader(input_path)
    writer = PyPDF2.PdfWriter()
    
    for page in reader.pages:
        page_width = float(page.mediabox.width)
        page_height = float(page.mediabox.height)
        
        crop_left = page_width * left
        crop_bottom = page_height * bottom
        crop_right = page_width * (1 - right)
        crop_top = page_height * (1 - top)
        
        page.cropbox.lower_left = (crop_left, crop_bottom)
        page.cropbox.upper_right = (crop_right, crop_top)
        
        writer.add_page(page)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

@app.route('/unlock-pdf', methods=['POST'])
def unlock_pdf():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        password = request.form.get('password', '')
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'unlocked_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        unlock_pdf_file(temp_input.name, output_path, password)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Unlock PDF failed: {str(e)}'}), 500

def unlock_pdf_file(input_path, output_path, password):
    reader = PyPDF2.PdfReader(input_path)
    
    if reader.is_encrypted:
        if not reader.decrypt(password):
            raise Exception('Invalid password')
    
    writer = PyPDF2.PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

@app.route('/protect-pdf', methods=['POST'])
def protect_pdf():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        password = request.form.get('password', '')
        if not password:
            return jsonify({'error': 'Password is required'}), 400
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'protected_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        protect_pdf_file(temp_input.name, output_path, password)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Protect PDF failed: {str(e)}'}), 500

def protect_pdf_file(input_path, output_path, password):
    reader = PyPDF2.PdfReader(input_path)
    writer = PyPDF2.PdfWriter()
    
    for page in reader.pages:
        writer.add_page(page)
    
    writer.encrypt(password)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

@app.route('/sign-pdf', methods=['POST'])
def sign_pdf():
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No PDF file provided'}), 400
            
        file = files[0]
        if not allowed_file(file.filename, 'pdf'):
            return jsonify({'error': 'Invalid PDF file'}), 400
        
        signature_text = request.form.get('signatureText', 'SIGNED')
        position = request.form.get('position', 'bottom-right')
        
        temp_input = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        file.save(temp_input.name)
        temp_input.close()
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'signed_pdf_{timestamp}.pdf'
        output_path = tempfile.mktemp(suffix='.pdf')
        
        sign_pdf_file(temp_input.name, output_path, signature_text, position)
        os.unlink(temp_input.name)
        
        return send_file_and_cleanup(output_path, output_filename)
    
    except Exception as e:
        return jsonify({'error': f'Sign PDF failed: {str(e)}'}), 500

def sign_pdf_file(input_path, output_path, signature_text, position):
    reader = PyPDF2.PdfReader(input_path)
    writer = PyPDF2.PdfWriter()
    
    for i, page in enumerate(reader.pages):
        if i == len(reader.pages) - 1:
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=letter)
            
            if position == 'bottom-left':
                x, y = 50, 50
            elif position == 'bottom-right':
                x, y = 400, 50
            elif position == 'top-left':
                x, y = 50, 750
            else:
                x, y = 400, 750
            
            can.setFont("Helvetica-Bold", 12)
            can.drawString(x, y, f"Digitally Signed: {signature_text}")
            can.drawString(x, y-15, f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
            can.save()
            
            packet.seek(0)
            overlay = PyPDF2.PdfReader(packet)
            page.merge_page(overlay.pages[0])
        
        writer.add_page(page)
    
    with open(output_path, 'wb') as output_file:
        writer.write(output_file)

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False, threaded=True)