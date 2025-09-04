#!/usr/bin/env python3
"""Test script to check PDF to PowerPoint dependencies"""

def test_ppt_dependencies():
    print("Testing PDF to PowerPoint dependencies...")
    
    try:
        import fitz
        print("✅ PyMuPDF (fitz) is installed")
    except ImportError:
        print("❌ PyMuPDF is NOT installed")
        print("Run: pip install PyMuPDF")
        return False
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        print("✅ python-pptx is installed")
        
        # Test creating a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        print("✅ PowerPoint creation test successful")
        
    except ImportError:
        print("❌ python-pptx is NOT installed")
        print("Run: pip install python-pptx")
        return False
    except Exception as e:
        print(f"❌ Error testing PowerPoint creation: {e}")
        return False
    
    print("✅ All PowerPoint dependencies are working!")
    return True

if __name__ == "__main__":
    test_ppt_dependencies()